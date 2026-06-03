from flask import Flask, request, jsonify, render_template, send_file
from flask_cors import CORS
import os, requests, json, re, tempfile, uuid
from datetime import datetime

app = Flask(__name__)
CORS(app)

# ═══════════════════════════════════════════════
# API KEYS (Render.com Environment Variables)
# ═══════════════════════════════════════════════
GROQ_KEY        = os.environ.get('GROQ_KEY')
GEMINI_KEY      = os.environ.get('GEMINI_KEY')
WEATHER_KEY     = os.environ.get('WEATHER_KEY')
SPOTIFY_ID      = os.environ.get('SPOTIFY_CLIENT_ID')
SPOTIFY_SECRET  = os.environ.get('SPOTIFY_CLIENT_SECRET')
YOUTUBE_KEY     = os.environ.get('YOUTUBE_KEY')
TELEGRAM_TOKEN  = os.environ.get('TELEGRAM_TOKEN')
TELEGRAM_CHAT   = os.environ.get('TELEGRAM_CHAT_ID')
NEWS_KEY        = os.environ.get('NEWS_KEY')
WOLFRAM_KEY     = os.environ.get('WOLFRAM_KEY')
NOTION_KEY      = os.environ.get('NOTION_KEY')
SERPAPI_KEY     = os.environ.get('SERPAPI_KEY')
OPENAI_KEY      = os.environ.get('OPENAI_KEY')
MAPS_KEY        = os.environ.get('MAPS_KEY')

# ═══════════════════════════════════════════════
# PING
# ═══════════════════════════════════════════════
@app.route('/ping')
def ping():
    return jsonify({'status': 'online', 'time': datetime.now().isoformat()})

@app.route('/')
def index():
    return render_template('index.html')

# ═══════════════════════════════════════════════
# CHAT — Groq (kısa) veya Gemini (uzun)
# ═══════════════════════════════════════════════
@app.route('/chat', methods=['POST'])
def chat():
    data = request.json
    message = data.get('message', '')
    long_mode = len(message) > 200 or any(k in message.lower() for k in ['anlat', 'açıkla', 'proje', 'ödev', 'makale', 'rapor'])

    system_prompt = """Sen J.A.R.V.I.S'sin — Iron Man filmindeki gibi gelişmiş bir yapay zeka asistanısın. 
Kullanıcına 'efendim' diye hitap et. Türkçe konuş. Kısa, net, zeki cevaplar ver.
Gerektiğinde komik ve espirili ol ama her zaman profesyonel kal.
Kullanıcının her isteğini yerine getirmeye çalış."""

    if long_mode and GEMINI_KEY:
        try:
            url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-1.5-flash:generateContent?key={GEMINI_KEY}"
            body = {"contents": [{"parts": [{"text": f"{system_prompt}\n\nKullanıcı: {message}"}]}]}
            r = requests.post(url, json=body, timeout=30)
            reply = r.json()['candidates'][0]['content']['parts'][0]['text']
            return jsonify({'reply': reply, 'engine': 'gemini'})
        except:
            pass

    if GROQ_KEY:
        try:
            headers = {'Authorization': f'Bearer {GROQ_KEY}', 'Content-Type': 'application/json'}
            body = {
                "model": "llama3-70b-8192",
                "messages": [
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": message}
                ],
                "max_tokens": 1024
            }
            r = requests.post('https://api.groq.com/openai/v1/chat/completions', headers=headers, json=body, timeout=20)
            resp_json = r.json()
            if 'choices' in resp_json:
                reply = resp_json['choices'][0]['message']['content']
                return jsonify({'reply': reply, 'engine': 'groq'})
            else:
                error_msg = resp_json.get('error', {}).get('message', str(resp_json))
                return jsonify({'reply': f'Groq hatası efendim: {error_msg}', 'engine': 'error'})
        except Exception as e:
            return jsonify({'reply': f'Bağlantı hatası efendim: {str(e)}', 'engine': 'error'})

    return jsonify({'reply': 'API anahtarı bulunamadı efendim.', 'engine': 'none'})

# ═══════════════════════════════════════════════
# WEATHER
# ═══════════════════════════════════════════════
@app.route('/weather')
def weather():
    city = request.args.get('city', 'Istanbul')
    try:
        url = f'https://api.openweathermap.org/data/2.5/weather?q={city}&appid={WEATHER_KEY}&units=metric&lang=tr'
        d = requests.get(url, timeout=10).json()
        return jsonify({
            'temp': round(d['main']['temp']),
            'feels_like': round(d['main']['feels_like']),
            'humidity': d['main']['humidity'],
            'desc': d['weather'][0]['description'].capitalize(),
            'wind_speed': round(d['wind']['speed'], 1),
            'city': d['name']
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# ═══════════════════════════════════════════════
# NEWS
# ═══════════════════════════════════════════════
@app.route('/news')
def news():
    category = request.args.get('category', 'general')
    lang = request.args.get('lang', 'tr')
    try:
        url = f'https://newsapi.org/v2/top-headlines?language={lang}&category={category}&pageSize=10&apiKey={NEWS_KEY}'
        d = requests.get(url, timeout=10).json()
        articles = []
        for a in d.get('articles', [])[:8]:
            articles.append({
                'title': a.get('title', ''),
                'source': a.get('source', {}).get('name', ''),
                'url': a.get('url', ''),
                'publishedAt': a.get('publishedAt', '')
            })
        return jsonify({'articles': articles})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# ═══════════════════════════════════════════════
# WOLFRAM ALPHA
# ═══════════════════════════════════════════════
@app.route('/wolfram')
def wolfram():
    query = request.args.get('q', '')
    try:
        url = f'https://api.wolframalpha.com/v1/result?i={requests.utils.quote(query)}&appid={WOLFRAM_KEY}'
        r = requests.get(url, timeout=10)
        return jsonify({'result': r.text})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# ═══════════════════════════════════════════════
# WIKIPEDIA
# ═══════════════════════════════════════════════
@app.route('/wiki')
def wiki():
    query = request.args.get('q', '')
    lang = request.args.get('lang', 'tr')
    try:
        url = f'https://{lang}.wikipedia.org/api/rest_v1/page/summary/{requests.utils.quote(query)}'
        d = requests.get(url, timeout=10).json()
        return jsonify({
            'title': d.get('title', ''),
            'summary': d.get('extract', ''),
            'url': d.get('content_urls', {}).get('desktop', {}).get('page', '')
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# ═══════════════════════════════════════════════
# WEB SEARCH (SerpAPI)
# ═══════════════════════════════════════════════
@app.route('/search')
def search():
    query = request.args.get('q', '')
    try:
        url = f'https://serpapi.com/search?q={requests.utils.quote(query)}&api_key={SERPAPI_KEY}&hl=tr&gl=tr'
        d = requests.get(url, timeout=15).json()
        results = []
        for r in d.get('organic_results', [])[:5]:
            results.append({
                'title': r.get('title', ''),
                'snippet': r.get('snippet', ''),
                'url': r.get('link', '')
            })
        return jsonify({'results': results})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# ═══════════════════════════════════════════════
# YOUTUBE SEARCH
# ═══════════════════════════════════════════════
@app.route('/youtube')
def youtube():
    query = request.args.get('q', '')
    try:
        url = f'https://www.googleapis.com/youtube/v3/search?part=snippet&q={requests.utils.quote(query)}&key={YOUTUBE_KEY}&maxResults=5&type=video'
        d = requests.get(url, timeout=10).json()
        videos = []
        for item in d.get('items', []):
            videos.append({
                'title': item['snippet']['title'],
                'channel': item['snippet']['channelTitle'],
                'videoId': item['id']['videoId'],
                'thumbnail': item['snippet']['thumbnails']['medium']['url'],
                'url': f"https://www.youtube.com/watch?v={item['id']['videoId']}"
            })
        return jsonify({'videos': videos})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# ═══════════════════════════════════════════════
# TELEGRAM — Bildirim Gönder
# ═══════════════════════════════════════════════
@app.route('/telegram/send', methods=['POST'])
def telegram_send():
    data = request.json
    message = data.get('message', '')
    try:
        url = f'https://api.telegram.org/bot{TELEGRAM_TOKEN}/sendMessage'
        r = requests.post(url, json={'chat_id': TELEGRAM_CHAT, 'text': message, 'parse_mode': 'HTML'}, timeout=10)
        return jsonify({'ok': r.json().get('ok', False)})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# ═══════════════════════════════════════════════
# SPOTIFY
# ═══════════════════════════════════════════════
spotify_token_cache = {'token': None, 'expires': 0}

def get_spotify_token():
    import time, base64
    if spotify_token_cache['token'] and time.time() < spotify_token_cache['expires']:
        return spotify_token_cache['token']
    creds = base64.b64encode(f"{SPOTIFY_ID}:{SPOTIFY_SECRET}".encode()).decode()
    r = requests.post('https://accounts.spotify.com/api/token',
        headers={'Authorization': f'Basic {creds}'},
        data={'grant_type': 'client_credentials'}, timeout=10)
    d = r.json()
    spotify_token_cache['token'] = d.get('access_token')
    spotify_token_cache['expires'] = time.time() + d.get('expires_in', 3600) - 60
    return spotify_token_cache['token']

@app.route('/spotify/search')
def spotify_search():
    query = request.args.get('q', '')
    try:
        token = get_spotify_token()
        url = f'https://api.spotify.com/v1/search?q={requests.utils.quote(query)}&type=track,artist&limit=5'
        d = requests.get(url, headers={'Authorization': f'Bearer {token}'}, timeout=10).json()
        tracks = []
        for t in d.get('tracks', {}).get('items', []):
            tracks.append({
                'name': t['name'],
                'artist': t['artists'][0]['name'],
                'album': t['album']['name'],
                'uri': t['uri'],
                'url': t['external_urls']['spotify']
            })
        return jsonify({'tracks': tracks})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/spotify/now-playing')
def spotify_now():
    # Bu endpoint kullanıcı OAuth token'ı gerektirir
    # Şimdilik placeholder
    return jsonify({'playing': False, 'message': 'OAuth gerekli'})

# ═══════════════════════════════════════════════
# GOOGLE MAPS — Rota / Yakın Yer
# ═══════════════════════════════════════════════
@app.route('/maps/directions')
def maps_directions():
    origin = request.args.get('from', '')
    destination = request.args.get('to', '')
    try:
        url = f'https://maps.googleapis.com/maps/api/directions/json?origin={requests.utils.quote(origin)}&destination={requests.utils.quote(destination)}&key={MAPS_KEY}&language=tr'
        d = requests.get(url, timeout=10).json()
        if d['routes']:
            leg = d['routes'][0]['legs'][0]
            return jsonify({
                'distance': leg['distance']['text'],
                'duration': leg['duration']['text'],
                'start': leg['start_address'],
                'end': leg['end_address']
            })
        return jsonify({'error': 'Rota bulunamadı'})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/maps/nearby')
def maps_nearby():
    location = request.args.get('location', '')
    place_type = request.args.get('type', 'restaurant')
    try:
        # Önce geocode
        geo_url = f'https://maps.googleapis.com/maps/api/geocode/json?address={requests.utils.quote(location)}&key={MAPS_KEY}'
        geo = requests.get(geo_url, timeout=10).json()
        if not geo['results']:
            return jsonify({'error': 'Konum bulunamadı'})
        loc = geo['results'][0]['geometry']['location']
        lat, lng = loc['lat'], loc['lng']

        nearby_url = f'https://maps.googleapis.com/maps/api/place/nearbysearch/json?location={lat},{lng}&radius=1500&type={place_type}&key={MAPS_KEY}&language=tr'
        d = requests.get(nearby_url, timeout=10).json()
        places = []
        for p in d.get('results', [])[:5]:
            places.append({
                'name': p.get('name'),
                'rating': p.get('rating'),
                'vicinity': p.get('vicinity')
            })
        return jsonify({'places': places})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# ═══════════════════════════════════════════════
# NOTION — Sayfa Oluştur / Listele
# ═══════════════════════════════════════════════
@app.route('/notion/pages')
def notion_pages():
    try:
        headers = {
            'Authorization': f'Bearer {NOTION_KEY}',
            'Notion-Version': '2022-06-28',
            'Content-Type': 'application/json'
        }
        r = requests.post('https://api.notion.com/v1/search',
            headers=headers, json={'filter': {'value': 'page', 'property': 'object'}}, timeout=10)
        results = r.json().get('results', [])
        pages = []
        for p in results[:10]:
            title = ''
            props = p.get('properties', {})
            for key in props:
                if props[key].get('type') == 'title':
                    texts = props[key].get('title', [])
                    if texts:
                        title = texts[0].get('plain_text', '')
            pages.append({'id': p['id'], 'title': title, 'url': p.get('url', '')})
        return jsonify({'pages': pages})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/notion/create', methods=['POST'])
def notion_create():
    data = request.json
    title = data.get('title', 'Yeni Sayfa')
    content = data.get('content', '')
    parent_id = data.get('parent_id', '')
    try:
        headers = {
            'Authorization': f'Bearer {NOTION_KEY}',
            'Notion-Version': '2022-06-28',
            'Content-Type': 'application/json'
        }
        body = {
            "parent": {"page_id": parent_id} if parent_id else {"type": "workspace", "workspace": True},
            "properties": {"title": {"title": [{"text": {"content": title}}]}},
            "children": [{"object": "block", "type": "paragraph",
                "paragraph": {"rich_text": [{"text": {"content": content}}]}}] if content else []
        }
        r = requests.post('https://api.notion.com/v1/pages', headers=headers, json=body, timeout=10)
        d = r.json()
        return jsonify({'id': d.get('id'), 'url': d.get('url')})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# ═══════════════════════════════════════════════
# DOSYA OLUŞTURMA — PDF / PPTX / DOCX
# ═══════════════════════════════════════════════
@app.route('/create-file', methods=['POST'])
def create_file_endpoint():
    data = request.json
    file_type = data.get('type', 'pdf')   # pdf | pptx | docx
    title = data.get('title', 'JARVIS Belgesi')
    content = data.get('content', '')
    topic = data.get('topic', '')

    # Önce SerpAPI ile araştır
    research = ''
    if topic and SERPAPI_KEY:
        try:
            url = f'https://serpapi.com/search?q={requests.utils.quote(topic)}&api_key={SERPAPI_KEY}&hl=tr'
            sr = requests.get(url, timeout=15).json()
            snippets = [r.get('snippet', '') for r in sr.get('organic_results', [])[:5]]
            research = '\n'.join(snippets)
        except:
            pass

    # Gemini ile içerik oluştur
    full_content = content
    if topic and GEMINI_KEY:
        try:
            prompt = f"""Sen JARVIS'sın. Aşağıdaki konu hakkında Türkçe detaylı bir belge içeriği oluştur.
Konu: {topic}
Araştırma verileri: {research}
Format: Başlık, giriş, ana bölümler (alt başlıklarla), sonuç şeklinde düzenle.
Belge türü: {file_type}"""
            url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-1.5-flash:generateContent?key={GEMINI_KEY}"
            body = {"contents": [{"parts": [{"text": prompt}]}]}
            r = requests.post(url, json=body, timeout=30)
            full_content = r.json()['candidates'][0]['content']['parts'][0]['text']
        except:
            full_content = content or f"{title} hakkında içerik oluşturulamadı."

    # Dosya oluştur
    tmp_dir = tempfile.gettempdir()
    filename = f"jarvis_{uuid.uuid4().hex[:8]}"

    try:
        if file_type == 'pdf':
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.units import cm
            from reportlab.pdfbase import pdfmetrics
            from reportlab.pdfbase.ttfonts import TTFont
            from reportlab.lib import colors

            filepath = os.path.join(tmp_dir, filename + '.pdf')
            doc = SimpleDocTemplate(filepath, pagesize=A4,
                rightMargin=2*cm, leftMargin=2*cm, topMargin=2*cm, bottomMargin=2*cm)
            styles = getSampleStyleSheet()
            story = []

            title_style = ParagraphStyle('Title', parent=styles['Title'],
                fontSize=18, spaceAfter=20, textColor=colors.HexColor('#003366'))
            body_style = ParagraphStyle('Body', parent=styles['Normal'],
                fontSize=11, spaceAfter=8, leading=16)

            story.append(Paragraph(title, title_style))
            story.append(Spacer(1, 0.5*cm))

            for line in full_content.split('\n'):
                line = line.strip()
                if not line:
                    story.append(Spacer(1, 0.3*cm))
                elif line.startswith('##'):
                    h_style = ParagraphStyle('H2', parent=styles['Heading2'],
                        fontSize=13, textColor=colors.HexColor('#0055aa'))
                    story.append(Paragraph(line.replace('##', '').strip(), h_style))
                elif line.startswith('#'):
                    h_style = ParagraphStyle('H1', parent=styles['Heading1'],
                        fontSize=15, textColor=colors.HexColor('#003366'))
                    story.append(Paragraph(line.replace('#', '').strip(), h_style))
                else:
                    story.append(Paragraph(line, body_style))

            doc.build(story)
            return send_file(filepath, as_attachment=True,
                download_name=f"{title}.pdf", mimetype='application/pdf')

        elif file_type == 'pptx':
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor

            filepath = os.path.join(tmp_dir, filename + '.pptx')
            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)

            # Slaytları böl
            sections = full_content.split('\n\n')
            slide_contents = []
            current_title = title
            current_body = []

            for section in sections:
                lines = section.strip().split('\n')
                if lines[0].startswith('#'):
                    if current_body:
                        slide_contents.append((current_title, '\n'.join(current_body)))
                        current_body = []
                    current_title = lines[0].replace('#', '').strip()
                    current_body = lines[1:]
                else:
                    current_body.extend(lines)

            if current_body:
                slide_contents.append((current_title, '\n'.join(current_body)))

            # Kapak slaytı
            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)
            bg = slide.background.fill
            bg.solid()
            bg.fore_color.rgb = RGBColor(0, 20, 60)

            txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.runs[0].font.size = Pt(40)
            p.runs[0].font.bold = True
            p.runs[0].font.color.rgb = RGBColor(0, 212, 255)

            # İçerik slaytları
            for slide_title, slide_body in slide_contents[:15]:
                slide = prs.slides.add_slide(blank_layout)
                bg = slide.background.fill
                bg.solid()
                bg.fore_color.rgb = RGBColor(0, 15, 45)

                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
                tf = title_box.text_frame
                p = tf.paragraphs[0]
                p.text = slide_title[:80]
                p.runs[0].font.size = Pt(24)
                p.runs[0].font.bold = True
                p.runs[0].font.color.rgb = RGBColor(0, 212, 255)

                body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5.5))
                tf = body_box.text_frame
                tf.word_wrap = True
                body_lines = slide_body.strip().split('\n')
                for i, line in enumerate(body_lines[:12]):
                    if not line.strip():
                        continue
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = line.strip()
                    p.runs[0].font.size = Pt(14)
                    p.runs[0].font.color.rgb = RGBColor(220, 235, 255)

            prs.save(filepath)
            return send_file(filepath, as_attachment=True,
                download_name=f"{title}.pptx",
                mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')

        elif file_type == 'docx':
            from docx import Document
            from docx.shared import Pt, RGBColor, Cm
            from docx.enum.text import WD_ALIGN_PARAGRAPH

            filepath = os.path.join(tmp_dir, filename + '.docx')
            doc = Document()

            # Sayfa kenar boşlukları
            section = doc.sections[0]
            section.top_margin = Cm(2.5)
            section.bottom_margin = Cm(2.5)
            section.left_margin = Cm(3)
            section.right_margin = Cm(3)

            # Başlık
            h = doc.add_heading(title, 0)
            h.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = h.runs[0]
            run.font.color.rgb = RGBColor(0, 51, 102)

            doc.add_paragraph()

            for line in full_content.split('\n'):
                line = line.strip()
                if not line:
                    doc.add_paragraph()
                elif line.startswith('## '):
                    doc.add_heading(line[3:], level=2)
                elif line.startswith('# '):
                    doc.add_heading(line[2:], level=1)
                elif line.startswith('- ') or line.startswith('* '):
                    p = doc.add_paragraph(line[2:], style='List Bullet')
                else:
                    p = doc.add_paragraph(line)
                    p.runs[0].font.size = Pt(11) if p.runs else None

            doc.save(filepath)
            return send_file(filepath, as_attachment=True,
                download_name=f"{title}.docx",
                mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document')

    except Exception as e:
        return jsonify({'error': f'Dosya oluşturma hatası: {str(e)}'}), 500

    return jsonify({'error': 'Geçersiz dosya türü'}), 400

# ═══════════════════════════════════════════════
# INTENT DETECTION — JARVIS komut algılama
# ═══════════════════════════════════════════════
@app.route('/intent', methods=['POST'])
def detect_intent():
    """
    JARVIS'ın mesajı analiz edip hangi servisi kullanacağını belirler.
    Örn: "spotify'da müzik çal" → {intent: 'spotify', action: 'search', query: 'müzik'}
    """
    data = request.json
    message = data.get('message', '').lower()

    intent = 'chat'
    params = {}

    # Spotify
    if any(k in message for k in ['müzik çal', 'şarkı çal', 'spotify', 'şarkı aç', 'müzik aç']):
        intent = 'spotify'
        params['action'] = 'search'
        params['query'] = message

    # YouTube
    elif any(k in message for k in ['youtube', 'video', 'izle']):
        intent = 'youtube'
        params['query'] = message

    # Haber
    elif any(k in message for k in ['haber', 'gündem', 'son dakika']):
        intent = 'news'

    # Hava durumu
    elif any(k in message for k in ['hava', 'sıcaklık', 'yağmur', 'kar']):
        intent = 'weather'

    # Dosya oluştur
    elif any(k in message for k in ['pdf', 'powerpoint', 'word', 'sunum', 'ödev', 'proje hazırla', 'rapor oluştur', 'belge']):
        intent = 'create_file'
        if 'pdf' in message:
            params['type'] = 'pdf'
        elif any(k in message for k in ['powerpoint', 'sunum', 'pptx']):
            params['type'] = 'pptx'
        else:
            params['type'] = 'docx'
        params['topic'] = message

    # Arama
    elif any(k in message for k in ['ara', 'bul', 'araştır', 'araştır', 'hakkında bilgi']):
        intent = 'search'
        params['query'] = message

    # Telegram
    elif any(k in message for k in ['telegram', 'bildirim gönder', 'mesaj gönder']):
        intent = 'telegram'
        params['message'] = message

    # Harita / Yol
    elif any(k in message for k in ['yol', 'rota', 'nasıl giderim', 'nerede', 'yakın']):
        intent = 'maps'
        params['query'] = message

    # Notion
    elif any(k in message for k in ['notion', 'not al', 'kaydet notion']):
        intent = 'notion'
        params['action'] = 'create'

    # Wikipedia
    elif any(k in message for k in ['vikipedi', 'wikipedia', 'nedir', 'kimdir']):
        intent = 'wiki'
        params['query'] = message

    # Wolfram
    elif any(k in message for k in ['hesapla', 'matematik', 'formül', 'kaç']):
        intent = 'wolfram'
        params['query'] = message

    return jsonify({'intent': intent, 'params': params})

if __name__ == '__main__':
    app.run(debug=False, host='0.0.0.0', port=5000)
